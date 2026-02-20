import faiss
import numpy as np
import os
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from groq import Groq
from rank_bm25 import BM25Okapi


embedding_model = None
index = None
chunks = None
bm25 = None


# -------- TEXT EXTRACTION --------
def extract_text_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


# -------- CHUNKING --------
def chunk_text(text, chunk_size=500, overlap=50):
    chunks = []
    start = 0

    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap

    return chunks


# -------- SETUP HYBRID RAG --------
def setup_rag(uploaded_file):

    global embedding_model, index, chunks, bm25

    embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

    text = extract_text_from_ppt(uploaded_file)
    chunks = chunk_text(text)

    # -------- VECTOR INDEX --------
    embeddings = embedding_model.encode(chunks)
    dimension = embeddings.shape[1]

    index = faiss.IndexFlatL2(dimension)
    index.add(np.array(embeddings))

    # -------- BM25 KEYWORD INDEX --------
    tokenized_chunks = [chunk.split() for chunk in chunks]
    bm25 = BM25Okapi(tokenized_chunks)


# -------- HYBRID SEARCH --------
def hybrid_search(query, top_k=3):

    # VECTOR SEARCH
    query_embedding = embedding_model.encode([query])
    distances, vector_indices = index.search(np.array(query_embedding), top_k)

    # KEYWORD SEARCH
    tokenized_query = query.split()
    keyword_scores = bm25.get_scores(tokenized_query)

    # Get top keyword indices
    keyword_indices = np.argsort(keyword_scores)[::-1][:top_k]

    # Combine results (union of both)
    combined_indices = list(set(vector_indices[0]) | set(keyword_indices))

    return [chunks[i] for i in combined_indices]


# -------- GET ANSWER --------
def get_answer(query):

    context_chunks = hybrid_search(query)
    context = "\n\n".join(context_chunks)

    client = Groq(api_key=os.getenv("GROQ_API_KEY"))

    prompt = f"""
    You are a professional AI assistant.

    Answer ONLY using the provided context.
    If answer not found say:
    "Information not available in document."

    Context:
    {context}

    Question:
    {query}
    """

    try:
        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[{"role": "user", "content": prompt}],
        )

        return response.choices[0].message.content, context_chunks

    except Exception as e:
        return f"Actual Error: {str(e)}", []